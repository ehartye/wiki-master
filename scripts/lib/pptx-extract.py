#!/usr/bin/env python3
"""Extract slide text/tables/notes from a .pptx into structured markdown on stdout.

Invoked by scripts/clip-pptx.mjs via execFileSync('python3', [thisFile, pptxPath]).
Self-contained: this file is the entire extraction technique, bundled inside
wiki-master, so the plugin has no runtime dependency on any other installed
plugin's copy of python-pptx glue. python-pptx itself is still an external,
independently-installable dependency (pip3 install python-pptx) -- exactly like
pandoc/soffice are for the sibling docx/xlsx clippers -- so any failure to import
it, or to open the given file, is reported on stderr with a non-zero exit so the
Node caller can classify it (missing dependency vs. a single bad file).
"""
import sys


def main():
    if len(sys.argv) < 2:
        print('usage: pptx-extract.py <file.pptx>', file=sys.stderr)
        sys.exit(2)

    try:
        from pptx import Presentation
    except ImportError:
        print('python-pptx is not installed (pip3 install python-pptx)', file=sys.stderr)
        sys.exit(1)

    path = sys.argv[1]
    try:
        prs = Presentation(path)
    except Exception as exc:  # corrupt/unreadable file -- fail clearly, never fabricate
        print(f'failed to open {path}: {exc}', file=sys.stderr)
        sys.exit(1)

    try:
        for i, slide in enumerate(prs.slides, start=1):
            # Emit the heading even for a slide with zero extractable content, so
            # slide numbering in the vault still matches the real deck -- useful
            # for citing "slide 7 says X" even when slide 7 is e.g. a title card.
            print(f'## Slide {i}')
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            print(t)
                if shape.has_table:
                    tbl = shape.table
                    n_cols = len(tbl.columns)
                    rows = [[tbl.cell(r, c).text.strip() for c in range(n_cols)] for r in range(len(tbl.rows))]
                    if rows:
                        header, *body = rows
                        print('| ' + ' | '.join(header) + ' |')
                        print('|' + '---|' * n_cols)
                        for row in body:
                            print('| ' + ' | '.join(row) + ' |')
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    print(f'**Speaker notes:** {notes}')
            print()
    except Exception as exc:  # a mid-deck failure is still "this file failed"
        print(f'failed to extract {path}: {exc}', file=sys.stderr)
        sys.exit(1)


if __name__ == '__main__':
    main()
